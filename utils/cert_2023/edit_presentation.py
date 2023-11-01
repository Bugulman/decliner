# %%
import os
import collections
import collections.abc
from pptx import Presentation
from pprint import pprint
import pandas as pd
from transliterate import translit


# %%

filename = 'participants4.csv'
csv_separator = '\t'
presentation_template = 'template_2023.pptx'
ready_folder_name = 'ready_certs'

if not os.path.exists(ready_folder_name):
    os.mkdir(ready_folder_name)

# %%

trainees = pd.read_csv(filename, sep=csv_separator, index_col=False)

for ind, person in trainees.iterrows():

    filling_dict = {
        '@USERNAME@': person['Name'],
        '@coursename@': person['course'],
        '@completed@': '',
        '@course_date@': person['date'],
        '@lector@': person['lector'],
        '@duration@': person['duration']
    }
    if person['M_or_F'] == 'f':
        filling_dict['@completed@'] = 'успешно освоила учебный курс на тему:'
    else:
        filling_dict['@completed@'] = 'успешно освоил учебный курс на тему:'

    prs = Presentation(presentation_template)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for parag in shape.text_frame.paragraphs:
                for run in parag.runs:
                    if filling_dict.get(run.text, False):
                        run.text = filling_dict[run.text]

    cert_filename_format = "{folder}/{numb:02d}_{name}.pptx"
    trans_name = translit(person['Name'], 'ru',
                          reversed=True).replace(' ', '_')
    cert_filename = cert_filename_format.format(
        numb=ind, name=trans_name, folder=ready_folder_name)
    prs.save(cert_filename)

print('Work complete!')
