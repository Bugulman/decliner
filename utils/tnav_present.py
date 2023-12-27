# %%
import os
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
# from transliterate import translit


# %%
dfadaf  
presentation_template = 'ГДМ_01_10.pptx'
ready_folder_name = 'ready_presentations'

if not os.path.exists(ready_folder_name):
    os.mkdir(ready_folder_name)

# %%

prs = Presentation(presentation_template)
path = Path(os.getcwd()) # %%
def get_pic_in_folder(path): pic_files = []
    for file in os.listdir(path):
        if file.endswith('.JPG'):
            pic_files.append(file)
    return pic_files
# %%


def create_slide_with_title(layout, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    title = slide.shapes.title
    title.text = title_text
    return slide

# %%


def add_picture_to_slide(slide, picture_path):
    pic = slide.shapes.add_picture(picture_path, top=Inches(2), left=Inches(4))
    return pic


# %%

pictures = get_pic_in_folder(path / 'reports')
pictures[0]

# path = path.joinpath(ready_folder_name)
# path.joinpath('Screenshots', pictures[0])
# %%
for picture in pictures:
    slide1=create_slide_with_title(16, picture.split('.')[0])
    add_picture_to_slide(slide1, str(path.joinpath('reports', picture)))

# %%
# slide1 = create_slide_with_title(16, '1 area')
# add_picture_to_slide(slide1, image)
prs.save(path.joinpath('ready_presentation.pptx'))
print('Work complete!')
